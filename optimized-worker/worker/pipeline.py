"""
Document extraction and storage download.

Uses PyMuPDF (fitz) for PDF extraction — significantly faster than pdfjs-dist.
Supports: PDF, DOCX, PPTX, MD, TXT, images (OCR).
"""
from __future__ import annotations

import io
import logging
import time
from dataclasses import dataclass, field
from typing import Optional

logger = logging.getLogger(__name__)


@dataclass
class ExtractedPage:
    page_number: int
    text: str
    is_image: bool = False


@dataclass
class ExtractionResult:
    full_text: str
    pages: list[ExtractedPage] = field(default_factory=list)
    timings: dict = field(default_factory=dict)


async def extract_pdf(buffer: bytes) -> ExtractionResult:
    """Extract text from PDF using PyMuPDF with per-page OCR fallback."""
    import fitz  # PyMuPDF
    from .ocr import ocr_pil_image
    from .config import get_settings
    from PIL import Image

    settings = get_settings()
    min_chars = settings.ocr_min_text_chars
    max_ocr = settings.ocr_max_pages

    pages: list[ExtractedPage] = []
    timings = {}

    t0 = time.perf_counter()
    try:
        doc = fitz.open(stream=buffer, filetype="pdf")
        timings["pdf_open_s"] = round(time.perf_counter() - t0, 4)

        ocr_count = 0
        t_text = time.perf_counter()

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text").strip()

            if len(text) < min_chars and ocr_count < max_ocr:
                # Render to image for OCR
                t_ocr = time.perf_counter()
                pix = page.get_pixmap(dpi=200)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                ocr_text = ocr_pil_image(img)
                ocr_count += 1

                if len(ocr_text) > len(text):
                    pages.append(ExtractedPage(
                        page_number=page_num + 1, text=ocr_text, is_image=True
                    ))
                    continue

            pages.append(ExtractedPage(
                page_number=page_num + 1, text=text, is_image=len(text) == 0
            ))

        timings["text_extraction_s"] = round(time.perf_counter() - t_text, 4)
        timings["ocr_pages"] = ocr_count
        doc.close()

    except Exception as e:
        logger.warning("PDF parsing failed, attempting whole-buffer OCR: %s", e)
        from .ocr import ocr_image_buffer
        ocr_text = ocr_image_buffer(buffer)
        if ocr_text:
            pages.append(ExtractedPage(page_number=1, text=ocr_text, is_image=True))

    full_text = "\n\n".join(p.text for p in pages)
    if not pages:
        pages = [ExtractedPage(page_number=1, text=full_text)]

    return ExtractionResult(full_text=full_text, pages=pages, timings=timings)


async def extract_docx(buffer: bytes) -> ExtractionResult:
    """Extract text from DOCX using mammoth."""
    import mammoth

    t0 = time.perf_counter()
    result = mammoth.extract_raw_text(io.BytesIO(buffer))
    text = (result.value or "").strip()
    elapsed = time.perf_counter() - t0

    return ExtractionResult(
        full_text=text,
        pages=[ExtractedPage(page_number=1, text=text)],
        timings={"docx_extraction_s": round(elapsed, 4)},
    )


async def extract_pptx(buffer: bytes) -> ExtractionResult:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation

    t0 = time.perf_counter()
    prs = Presentation(io.BytesIO(buffer))
    texts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text_parts.append(paragraph.text)
        texts.append("\n".join(slide_text_parts))

    full_text = "\n\n".join(texts)
    elapsed = time.perf_counter() - t0

    return ExtractionResult(
        full_text=full_text,
        pages=[ExtractedPage(page_number=1, text=full_text)],
        timings={"pptx_extraction_s": round(elapsed, 4)},
    )


async def extract_markdown(buffer: bytes) -> ExtractionResult:
    """Extract text from Markdown, stripping YAML frontmatter."""
    import re

    text = buffer.decode("utf-8", errors="replace")
    text = re.sub(r"^---\s*\n[\s\S]*?\n---\s*\n", "", text).strip()
    return ExtractionResult(
        full_text=text,
        pages=[ExtractedPage(page_number=1, text=text)],
    )


async def extract_plain_text(buffer: bytes) -> ExtractionResult:
    """Extract plain text."""
    text = buffer.decode("utf-8", errors="replace").strip()
    return ExtractionResult(
        full_text=text,
        pages=[ExtractedPage(page_number=1, text=text)],
    )


async def extract_image(buffer: bytes) -> ExtractionResult:
    """OCR an image."""
    from .ocr import ocr_image_buffer

    t0 = time.perf_counter()
    text = ocr_image_buffer(buffer)
    elapsed = time.perf_counter() - t0

    return ExtractionResult(
        full_text=text,
        pages=[ExtractedPage(page_number=1, text=text, is_image=True)],
        timings={"ocr_s": round(elapsed, 4)},
    )


async def extract_document(
    file_name: str, buffer: bytes, mime_type: str = ""
) -> ExtractionResult:
    """Route extraction based on file extension/MIME type."""
    ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""
    mime = (mime_type or "").lower()

    if mime.startswith("application/pdf") or ext == "pdf":
        return await extract_pdf(buffer)
    if "presentation" in mime or ext == "pptx":
        return await extract_pptx(buffer)
    if "wordprocessingml" in mime or "msword" in mime or ext in ("docx", "doc"):
        return await extract_docx(buffer)
    if mime.startswith("image/") or ext in ("png", "jpg", "jpeg", "webp", "gif", "bmp", "tiff"):
        return await extract_image(buffer)
    if ext in ("md", "markdown") or "markdown" in mime:
        return await extract_markdown(buffer)

    return await extract_plain_text(buffer)
